import shutil
import os

from openpyxl import load_workbook
from odf.opendocument import load
from odf import text, teletype
from pptx import Presentation
from docx import Document
import PyPDF2

from app import tmp_path


# Reads the content from all the files in the provided folderS
def read_text(file_path, word_limit):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
        content = file.read().split()
        return ' '.join(content[:word_limit])

def read_pdf(file_path, word_limit):
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            content = ' '.join(page.extract_text() for page in pdf_reader.pages).split()
            return ' '.join(content[:word_limit])

def read_odf(file_path, word_limit):
        odf_file = load(file_path)
        content = [teletype.extractText(para) for para in odf_file.getElementsByType(text.P)]
        return ' '.join(content[:word_limit])

def read_doc(file_path, word_limit):
        doc = Document(file_path)
        content = ' '.join([paragraph.text for paragraph in doc.paragraphs]).split()
        return ' '.join(content[:word_limit])

def read_xlsx(file_path, word_limit):
        workbook = load_workbook(file_path)
        sheet = workbook.active
        rows = list(sheet.iter_rows(values_only=True))
        content = [cell for row in rows for cell in row if cell]
        return ' '.join(str(cell) for cell in content[:word_limit])

def read_ppt(file_path, word_limit):
        presentation = Presentation(file_path)
        content = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content.append(shape.text)
        return ' '.join(content[:word_limit])


# Avoids several ifs by mapping funcs to exts
func_map = {'.txt':  read_text, '.html': read_text, '.md': read_text, '.csv': read_text, 
            '.pdf':  read_pdf,  '.docx': read_doc,  '.odt': read_odf, '.odp': read_odf, 
            '.xlsx': read_xlsx, '.pptx': read_ppt,  '.ppt': read_ppt  } 

def prep_files(directory, select_folder, copy_files=False):
    # if user selects a folder: Moves the files from existing sub-folders (if any) to root level of selected folder
    if select_folder:
        for root, _, files in os.walk(directory):
            for file_name in files:
                file_path = os.path.join(root, file_name)
                root_path = os.path.join(directory, file_name)
                shutil.move(file_path, root_path)

        # Deletes the now empty existing sub-folders
        for root, folders, _ in os.walk(directory, topdown=False):
            for folder in folders:
                os.rmdir(os.path.join(root, folder))

    # If user uploads files: Moves/Copies the said files to tmp folder in the installation directory
    if not select_folder:
        tmp_folder = os.path.join(tmp_path, "Organized_Files")
        if not os.path.exists(tmp_folder):
            os.mkdir(tmp_folder)

        for file_path in directory:
            if copy_files:
                shutil.copy(file_path, os.path.join(tmp_folder, os.path.basename(file_path)))
            else:
                shutil.move(file_path, os.path.join(tmp_folder, os.path.basename(file_path)))

def read_files(directory, word_limit):
    words_list = []
    misc_list = []
    for _, _, files in os.walk(directory):
        for file_name in files:
            root_path = os.path.join(directory, file_name)

            # Reads the files and adds the data to words_list
            if os.path.isfile(root_path):
                file_extension = os.path.splitext(file_name)[1]
                if file_extension in func_map:
                    words_list.append((file_name, func_map[file_extension](root_path, word_limit)))
                else:
                     misc_list.append(file_name)

    return words_list, misc_list
