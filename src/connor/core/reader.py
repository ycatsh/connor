import shutil
import os
from pathlib import Path
from typing import (
    Callable, List, Tuple, Dict
)

from openpyxl import load_workbook
from odf.opendocument import load as load_odf
from odf import text, teletype
from pptx import Presentation
from docx import Document
import PyPDF2


def read_text(file_path: Path, word_limit: int) -> str:
    try:
        with file_path.open('r', encoding='utf-8', errors='ignore') as file:
            content = file.read().split()
            return ' '.join(content[:word_limit])
    except Exception as e:
        return f"Error reading {file_path.name}: {e}"


def read_pdf(file_path: Path, word_limit: int) -> str:
    try:
        with file_path.open('rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            content = []
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    content.extend(text.split())
            return ' '.join(content[:word_limit])
    except Exception as e:
        return f"Error reading {file_path.name}: {e}"


def read_odf(file_path: Path, word_limit: int) -> str:
    try:
        odf_file = load_odf(file_path)
        paras = odf_file.getElementsByType(text.P)
        content = []
        for para in paras:
            text_data = teletype.extractText(para)
            if text_data:
                content.extend(text_data.split())
        return ' '.join(content[:word_limit])
    except Exception as e:
        return f"Error reading {file_path.name}: {e}"


def read_doc(file_path: Path, word_limit: int) -> str:
    try:
        doc = Document(file_path)
        content = []
        for para in doc.paragraphs:
            content.extend(para.text.split())
        return ' '.join(content[:word_limit])
    except Exception as e:
        return f"Error reading {file_path.name}: {e}"


def read_xlsx(file_path: Path, word_limit: int) -> str:
    try:
        workbook = load_workbook(file_path)
        sheet = workbook.active
        content = []
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    content.extend(str(cell).split())
        return ' '.join(content[:word_limit])
    except Exception as e:
        return f"Error reading {file_path.name}: {e}"


def read_ppt(file_path: Path, word_limit: int) -> str:
    try:
        presentation = Presentation(file_path)
        content = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content.extend(shape.text.split())
        return ' '.join(content[:word_limit])
    except Exception as e:
        return f"Error reading {file_path.name}: {e}"


FUNC_MAP: Dict[str, Callable[[Path, int], str]] = {
    '.txt': read_text,
    '.html': read_text,
    '.md': read_text,
    '.csv': read_text,
    '.pdf': read_pdf,
    '.docx': read_doc,
    '.odt': read_odf,
    '.odp': read_odf,
    '.xlsx': read_xlsx,
    '.pptx': read_ppt,
    '.ppt': read_ppt
}


def prep_files(directory: Path) -> None:
    """
    Flattens the directory.
    """
    for root, _, files in os.walk(directory):
        root_path = Path(root)
        for file_name in files:
            file_path = root_path / file_name
            destination = directory / file_name
            if file_path != destination:
                shutil.move(str(file_path), str(destination))

    for root, folders, _ in os.walk(directory, topdown=False):
        root_path = Path(root)
        for folder in folders:
            folder_path = root_path / folder
            try:
                folder_path.rmdir()
            except OSError:
                print(f"Could not delete folder {folder_path}")


def read_files(folder_path: Path, word_limit: int) -> Tuple[List[Tuple[str, str]], List[str]]:
    """
    Reads files in the folder by reading words till the word limit.

    Args:
        folder_path: Path object of the selected folder to read.
        word_limit: Maximum words to read in each file.

    Returns:
        Tuple containing:
            List of tuples (filename, content) for processed text files
            List of filenames for miscellaneous/unreadable files
    """
    text_files_list = []
    misc_files_list = []
    for file_path in folder_path.iterdir():
        if file_path.is_file():
            extension = file_path.suffix.lower()
            if extension in FUNC_MAP:
                reader = FUNC_MAP[extension]
                content = reader(file_path, word_limit)
                text_files_list.append((file_path.name, content))
            else:
                misc_files_list.append(file_path.name)
    return text_files_list, misc_files_list
